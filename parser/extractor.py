"""
Extractor images from pptx and tag them
"""
import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.exc import PackageNotFoundError
import logging

logger = logging.getLogger(__name__)

class Extractor:
    """
    Extractor class
    """
    def __init__(self, output_dir : str):
        self.output_dir = output_dir
        os.makedirs(self.output_dir, exist_ok=True)

    def extract_images_from_ppt(self, ppt_path : str) -> None:
        """
        Extract images from a single PPT file and save them to the output directory
        """
        try:
            presentation = Presentation(ppt_path)
        except PackageNotFoundError as e:
            logger.error("Failed to open PPT file - %s [%s]", ppt_path, e)
            raise e

        ppt_name = os.path.splitext(os.path.basename(ppt_path))[0]
        for slide_index, slide in enumerate(presentation.slides):
            for shape_index, shape in enumerate(slide.shapes):
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image = shape.image
                    image_bytes = image.blob
                    image_format = image.ext

                    image_name = f"{ppt_name}_slide{slide_index}_img{shape_index}.{image_format}"
                    image_path = os.path.join(self.output_dir, image_name)

                    with open(image_path, 'wb') as f:
                        f.write(image_bytes)

    def find_pptx_files(self, input_folder : str) -> list[str]:
        """
        Find all PPTX files in the input folder
        """
        if not os.path.exists(input_folder):
            logger.error("Input folder %s does not exist", input_folder)
            raise ValueError(f"Input folder {input_folder} does not exist")

        ppt_files = []
        for filename in os.listdir(input_folder):
            if filename.lower().endswith('.pptx'):
                ppt_path = os.path.join(input_folder, filename)
                if os.path.isfile(ppt_path):
                    ppt_files.append(ppt_path)

        if not ppt_files:
            logger.warning("No PPTX files found in %s", input_folder)
            return []

        logger.info("Found %s PPTX files in %s", len(ppt_files), input_folder)
        return ppt_files
